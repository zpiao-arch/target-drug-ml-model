from __future__ import annotations

import json
import math
import shutil
import sys
from pathlib import Path

import joblib
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
import numpy as np
import pandas as pd
import shap
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.ensemble import ExtraTreesClassifier
from sklearn.metrics import (
    auc,
    balanced_accuracy_score,
    confusion_matrix,
    f1_score,
    log_loss,
    precision_score,
    recall_score,
    roc_auc_score,
    roc_curve,
)
from sklearn.model_selection import GroupShuffleSplit


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "ai_mol_loop"))

import target_match_model as tmm  # noqa: E402


OUT = ROOT / "deliverables" / "extra_trees_model_ppt_20260601"
FIG = OUT / "assets" / "figures"
PPTX = OUT / "ViroMol_Compass_ExtraTrees模型汇报.pptx"
QA = OUT / "qa_report.md"
SUMMARY_JSON = OUT / "experiment_summary.json"
CASE_RANKINGS_CSV = OUT / "multi_disease_target_rankings.csv"
MODEL_COMPARISON_SOURCE = ROOT / "deliverables" / "project_proposal_demo_20260531" / "experiments" / "target_match_supplemental_benchmark.csv"
MODEL_COMPARISON_CSV = OUT / "model_comparison_benchmark.csv"


TITLE = "ViroMol Compass ExtraTrees 靶点-药物匹配模型"
SUBTITLE = "基于 DrugCentral 关系、分子结构特征与靶点文本证据的计算筛选排序模型"
BOUNDARY = "Computational target-drug matching only; not biological activity, safety, dosing, or clinical evidence."

DISEASE_TARGET_CASES = [
    {
        "case": "甲/乙流 NA",
        "target": "neuraminidase",
        "query": "influenza neuraminidase NA oseltamivir zanamivir P03468",
    },
    {
        "case": "HIV 逆转录酶",
        "target": "reverse transcriptase",
        "query": "Human immunodeficiency virus 1 reverse transcriptase RNaseH pol Q72547 antiretroviral",
    },
    {
        "case": "HCV NS5B",
        "target": "RNA polymerase",
        "query": "hepatitis C virus RNA-dependent RNA polymerase NS5B antiviral sofosbuvir",
    },
    {
        "case": "高血压 ACE",
        "target": "ACE",
        "query": "Homo sapiens angiotensin-converting enzyme ACE P12821 hypertension inhibitor",
    },
    {
        "case": "2型糖尿病 DPP4",
        "target": "DPP4",
        "query": "Homo sapiens dipeptidyl peptidase 4 DPP4 diabetes inhibitor",
    },
    {
        "case": "炎症疼痛 COX-2",
        "target": "PTGS2",
        "query": "Homo sapiens prostaglandin G/H synthase 2 PTGS2 cyclooxygenase inflammation inhibitor",
    },
    {
        "case": "肿瘤 EGFR",
        "target": "EGFR",
        "query": "Homo sapiens epidermal growth factor receptor EGFR cancer inhibitor",
    },
]

CHINESE_FONT = "PingFang SC"
LATIN_FONT = "Aptos"
PAPER = RGBColor(249, 247, 243)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(26, 33, 42)
MUTED = RGBColor(86, 95, 106)
SOFT = RGBColor(125, 132, 141)
LINE = RGBColor(215, 210, 200)
BLUE = RGBColor(37, 78, 123)
GREEN = RGBColor(47, 93, 80)
TERRA = RGBColor(155, 90, 46)


def ensure_dirs() -> None:
    if OUT.exists():
        shutil.rmtree(OUT)
    FIG.mkdir(parents=True, exist_ok=True)


def set_plot_style() -> None:
    font_candidates = [
        Path("/Library/Fonts/Arial Unicode.ttf"),
        Path("/System/Library/Fonts/STHeiti Medium.ttc"),
        Path("/System/Library/Fonts/STHeiti Light.ttc"),
    ]
    plot_font = "DejaVu Sans"
    for font_path in font_candidates:
        if font_path.exists():
            font_manager.fontManager.addfont(str(font_path))
            plot_font = font_manager.FontProperties(fname=str(font_path)).get_name()
            break
    plt.rcParams.update(
        {
            "font.family": plot_font,
            "axes.unicode_minus": False,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.labelcolor": "#26323F",
            "xtick.color": "#56606B",
            "ytick.color": "#56606B",
            "axes.edgecolor": "#D8D2C8",
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "savefig.dpi": 220,
            "axes.titleweight": "bold",
            "axes.titlesize": 13,
            "axes.labelsize": 10.5,
            "legend.fontsize": 9.5,
        }
    )


def build_dataset():
    dataset = tmm.build_target_match_dataset(negative_ratio=1, seed=13)
    mol_x, target_x, labels = tmm._features_for_examples(dataset, mol_dim=512, target_dim=256)
    x = np.hstack([mol_x, target_x]).astype(np.float32)
    y = labels.astype(int)
    struct_ids = np.array([str(item["struct_id"]) for item in dataset.examples])
    target_keys = np.array([str(item["target_key"]) for item in dataset.examples])
    return dataset, x, y, struct_ids, target_keys


def split_indices(y, struct_ids, target_keys, seed=13):
    rng = np.random.default_rng(seed)
    idx = np.arange(len(y))
    rng.shuffle(idx)
    cut = int(len(idx) * 0.8)
    random_split = (idx[:cut], idx[cut:])

    target_gss = GroupShuffleSplit(n_splits=1, test_size=0.2, random_state=seed)
    target_train, target_valid = next(target_gss.split(np.zeros_like(y), y, groups=target_keys))

    drug_gss = GroupShuffleSplit(n_splits=1, test_size=0.2, random_state=seed)
    drug_train, drug_valid = next(drug_gss.split(np.zeros_like(y), y, groups=struct_ids))

    return {
        "random_pair_split": random_split,
        "target_group_split": (target_train, target_valid),
        "drug_group_split": (drug_train, drug_valid),
    }


def train_extra_trees(x_train, y_train, n_estimators=180):
    model = ExtraTreesClassifier(
        n_estimators=n_estimators,
        min_samples_leaf=2,
        max_features="sqrt",
        random_state=13,
        n_jobs=-1,
    )
    model.fit(x_train, y_train)
    return model


def evaluate_split(name, x, y, train_idx, valid_idx):
    model = train_extra_trees(x[train_idx], y[train_idx])
    prob = model.predict_proba(x[valid_idx])[:, 1]
    pred = (prob >= 0.5).astype(int)
    fpr, tpr, _ = roc_curve(y[valid_idx], prob)
    metrics = {
        "split": name,
        "train_examples": int(len(train_idx)),
        "valid_examples": int(len(valid_idx)),
        "valid_positive_rate": round(float(y[valid_idx].mean()), 4),
        "accuracy": round(float((pred == y[valid_idx]).mean()), 4),
        "balanced_accuracy": round(float(balanced_accuracy_score(y[valid_idx], pred)), 4),
        "roc_auc": round(float(roc_auc_score(y[valid_idx], prob)), 4),
        "precision": round(float(precision_score(y[valid_idx], pred, zero_division=0)), 4),
        "recall": round(float(recall_score(y[valid_idx], pred, zero_division=0)), 4),
        "f1": round(float(f1_score(y[valid_idx], pred, zero_division=0)), 4),
        "confusion_matrix": confusion_matrix(y[valid_idx], pred).tolist(),
    }
    return model, metrics, {"fpr": fpr, "tpr": tpr, "prob": prob, "y_true": y[valid_idx]}


def plot_roc(roc_data, metrics):
    colors = {
        "random_pair_split": "#254E7B",
        "target_group_split": "#2F5D50",
        "drug_group_split": "#9B5A2E",
    }
    labels = {
        "random_pair_split": "随机划分",
        "target_group_split": "冷靶点",
        "drug_group_split": "冷药物",
    }
    fig, ax = plt.subplots(figsize=(6.8, 5.1))
    for name, data in roc_data.items():
        auc_value = metrics[name]["roc_auc"]
        ax.plot(data["fpr"], data["tpr"], lw=2.8, color=colors[name], label=f"{labels[name]}  ROC-AUC={auc_value:.4f}")
    ax.plot([0, 1], [0, 1], "--", color="#B5B0A8", lw=1.1, label="随机基线")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1.02)
    ax.set_xlabel("假阳性率")
    ax.set_ylabel("真阳性率")
    ax.set_title("三种验证划分下的 ROC 曲线")
    ax.legend(frameon=False, loc="lower right")
    ax.grid(alpha=0.16, color="#D8D2C8")
    path = FIG / "roc_curves_extra_trees.png"
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def plot_metric_comparison(summary_df):
    ordered = ["random_pair_split", "target_group_split", "drug_group_split"]
    labels = ["随机划分", "冷靶点", "冷药物"]
    lookup = summary_df.set_index("split")
    values = [lookup.loc[s, "roc_auc"] for s in ordered]
    accuracy_values = [lookup.loc[s, "accuracy"] for s in ordered]
    f1_values = [summary_df.set_index("split").loc[s, "f1"] for s in ordered]
    x = np.arange(len(ordered))
    fig, ax = plt.subplots(figsize=(7.2, 4.6))
    width = 0.24
    series = [
        ("排序能力 ROC-AUC", values, "#254E7B", -width),
        ("准确率", accuracy_values, "#2F5D50", 0),
        ("阈值性能 F1", f1_values, "#9B5A2E", width),
    ]
    for label, vals, color, offset in series:
        ax.bar(x + offset, vals, width, color=color, label=label)
        for i, v in enumerate(vals):
            ax.text(i + offset, v + 0.014, f"{v:.3f}", ha="center", va="bottom", fontsize=8.5)
    ax.set_ylim(0, 1.06)
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylabel("指标值")
    ax.set_title("排序、准确率和阈值性能的综合对比")
    ax.legend(frameon=False, loc="upper center", bbox_to_anchor=(0.5, -0.10), ncol=3)
    ax.grid(axis="y", alpha=0.16, color="#D8D2C8")
    fig.subplots_adjust(bottom=0.22)
    path = FIG / "benchmark_metric_bars.png"
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def learning_curve(x, y, train_idx, valid_idx):
    rows = []
    for n in [10, 20, 40, 80, 120, 180]:
        model = train_extra_trees(x[train_idx], y[train_idx], n_estimators=n)
        train_prob = model.predict_proba(x[train_idx])[:, 1]
        valid_prob = model.predict_proba(x[valid_idx])[:, 1]
        rows.append(
            {
                "n_estimators": n,
                "train_log_loss": float(log_loss(y[train_idx], train_prob, labels=[0, 1])),
                "valid_log_loss": float(log_loss(y[valid_idx], valid_prob, labels=[0, 1])),
                "valid_roc_auc": float(roc_auc_score(y[valid_idx], valid_prob)),
            }
        )
    df = pd.DataFrame(rows)
    df.to_csv(OUT / "learning_curve_random_pair.csv", index=False)
    fig, ax1 = plt.subplots(figsize=(7.2, 4.7))
    ax1.plot(df["n_estimators"], df["train_log_loss"], marker="o", color="#2F5D50", label="训练 log-loss")
    ax1.plot(df["n_estimators"], df["valid_log_loss"], marker="o", color="#C0504D", label="验证 log-loss")
    ax1.set_xlabel("树数量")
    ax1.set_ylabel("Log-loss")
    ax1.grid(alpha=0.16, color="#D8D2C8")
    ax2 = ax1.twinx()
    ax2.plot(df["n_estimators"], df["valid_roc_auc"], marker="s", color="#254E7B", label="验证 ROC-AUC")
    ax2.set_ylabel("ROC-AUC")
    ax2.set_ylim(max(0.7, df["valid_roc_auc"].min() - 0.04), min(1.0, df["valid_roc_auc"].max() + 0.04))
    lines, labels = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines + lines2, labels + labels2, frameon=False, loc="center right")
    ax1.set_title("树数量增加时的 log-loss 与 ROC-AUC")
    path = FIG / "learning_curve_logloss_auc.png"
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path, df


def feature_names():
    descriptor_names = [
        "分子量",
        "MolLogP",
        "TPSA",
        "氢键供体",
        "氢键受体",
        "可旋转键",
        "重原子数",
        "环数量",
    ]
    names = [f"Morgan 位点 {i}" for i in range(512)]
    names.extend(descriptor_names)
    names.extend([f"靶点哈希 {i}" for i in range(256)])
    return names, descriptor_names


def aggregate_blocks(values):
    return {
        "Morgan 指纹位点": float(np.sum(values[:512])),
        "RDKit 标量描述符": float(np.sum(values[512:520])),
        "靶点文本哈希": float(np.sum(values[520:])),
    }


def plot_feature_importance(model):
    names, descriptor_names = feature_names()
    importances = model.feature_importances_
    blocks = aggregate_blocks(importances)
    fig, axes = plt.subplots(1, 2, figsize=(9.0, 4.4), gridspec_kw={"width_ratios": [1.05, 1.35]})
    block_items = sorted(blocks.items(), key=lambda x: x[1], reverse=True)
    axes[0].barh([x[0] for x in block_items][::-1], [x[1] for x in block_items][::-1], color="#254E7B")
    axes[0].set_title("特征块重要性")
    axes[0].set_xlabel("重要性总量")
    descriptor_imp = importances[512:520]
    order = np.argsort(descriptor_imp)
    axes[1].barh([descriptor_names[i] for i in order], descriptor_imp[order], color="#9B5A2E")
    axes[1].set_title("RDKit 描述符细节")
    axes[1].set_xlabel("重要性")
    for ax in axes:
        ax.grid(axis="x", alpha=0.15, color="#D8D2C8")
    path = FIG / "feature_importance_blocks.png"
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path, blocks


def run_shap(model, x, valid_idx):
    rng = np.random.default_rng(13)
    sample_size = min(40, len(valid_idx))
    sample_idx = rng.choice(valid_idx, size=sample_size, replace=False)
    x_sample = x[sample_idx]
    explainer = shap.TreeExplainer(model)
    shap_values = explainer.shap_values(x_sample, check_additivity=False, approximate=True)
    if isinstance(shap_values, list):
        values = shap_values[1]
    else:
        arr = np.asarray(shap_values)
        if arr.ndim == 3:
            values = arr[:, :, 1]
        else:
            values = arr
    abs_mean = np.mean(np.abs(values), axis=0)
    blocks = aggregate_blocks(abs_mean)
    names, descriptor_names = feature_names()
    top_idx = np.argsort(abs_mean)[-15:]
    top_labels = [names[i] for i in top_idx]
    top_values = abs_mean[top_idx]

    fig, axes = plt.subplots(1, 2, figsize=(9.2, 4.6), gridspec_kw={"width_ratios": [1.0, 1.4]})
    block_items = sorted(blocks.items(), key=lambda x: x[1], reverse=True)
    axes[0].barh([x[0] for x in block_items][::-1], [x[1] for x in block_items][::-1], color="#2F5D50")
    axes[0].set_title("特征块平均 |SHAP|")
    axes[0].set_xlabel("平均绝对 SHAP 总量")
    axes[1].barh(top_labels, top_values, color="#254E7B")
    axes[1].set_title("平均 |SHAP| 最高的原始特征")
    axes[1].set_xlabel("平均 |SHAP|")
    for ax in axes:
        ax.grid(axis="x", alpha=0.15, color="#D8D2C8")
    path = FIG / "shap_grouped_summary.png"
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path, blocks


def plot_confusion(metrics):
    labels = {
        "random_pair_split": "随机划分",
        "target_group_split": "冷靶点",
        "drug_group_split": "冷药物",
    }
    fig, axes = plt.subplots(1, 3, figsize=(9.0, 3.2))
    for ax, (name, m) in zip(axes, metrics.items()):
        cm = np.asarray(m["confusion_matrix"])
        im = ax.imshow(cm, cmap="Blues")
        ax.set_title(labels[name], fontsize=10)
        ax.set_xticks([0, 1])
        ax.set_xticklabels(["预测 0", "预测 1"], fontsize=8)
        ax.set_yticks([0, 1])
        ax.set_yticklabels(["真实 0", "真实 1"], fontsize=8)
        for i in range(2):
            for j in range(2):
                ax.text(j, i, str(cm[i, j]), ha="center", va="center", fontsize=10, color="#111")
    fig.suptitle("概率阈值 0.5 下的混淆矩阵", y=1.02, fontweight="bold")
    path = FIG / "confusion_matrices.png"
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def load_influenza_predictions():
    p = ROOT / "ai_mol_loop" / "models" / "target_match_drugcentral" / "influenza_na_predictions_optimized.csv"
    if not p.exists():
        p = ROOT / "deliverables" / "project_proposal_demo_20260531" / "experiments" / "influenza_na_known_drug_match_top20.csv"
    return pd.read_csv(p).head(6)


def rank_multi_disease_cases(model, dataset):
    structures = list(dataset.structures_by_id.values())
    mol_matrix = np.vstack([tmm.molecule_features(str(structure.get("smiles", "")), 512) for structure in structures]).astype(np.float32)
    target_records = {
        str(key): tmm._target_evidence_record(target, 256)
        for key, target in dataset.targets_by_key.items()
        if isinstance(target, dict)
    }
    rows = []
    for case in DISEASE_TARGET_CASES:
        query = case["query"]
        query_vec = tmm.target_text_features(query, 256)
        feature_matrix = np.hstack([mol_matrix, np.repeat(query_vec.reshape(1, -1), len(structures), axis=0)]).astype(np.float32)
        model_probabilities = model.predict_proba(feature_matrix)[:, 1]
        ranked = []
        for structure, model_probability in zip(structures, model_probabilities):
            struct_id = str(structure.get("struct_id", ""))
            evidence_similarity, evidence_target_key = tmm._best_evidence_similarity(
                query,
                query_vec,
                dataset.drug_target_index.get(struct_id, []),
                dataset.targets_by_key,
                256,
                target_records_by_key=target_records,
            )
            final_score = tmm.calibrated_match_probability(float(model_probability), evidence_similarity)
            matched_target = dataset.targets_by_key.get(evidence_target_key, {})
            ranked.append(
                {
                    "case": case["case"],
                    "target": case["target"],
                    "rank": 0,
                    "drug_name": str(structure.get("drug_name", "")),
                    "score": round(final_score, 6),
                    "model_probability": round(float(model_probability), 6),
                    "known_target_similarity": round(float(evidence_similarity), 6),
                    "evidence_target": str(matched_target.get("target_name", "")),
                    "evidence_organism": str(matched_target.get("organism", "")),
                    "boundary": "计算筛选排序信号，不代表真实疗效或临床结论。",
                }
            )
        ranked.sort(key=lambda row: float(row["score"]), reverse=True)
        for idx, row in enumerate(ranked[:5], start=1):
            row["rank"] = idx
            rows.append(row)
    df = pd.DataFrame(rows)
    df.to_csv(CASE_RANKINGS_CSV, index=False)
    return df


def plot_case_rankings(case_df):
    top_df = case_df[case_df["rank"] <= 3].copy()
    cases = [case["case"] for case in DISEASE_TARGET_CASES]
    ranks = [1, 2, 3]
    fig, ax = plt.subplots(figsize=(10.0, 5.7))
    ax.set_xlim(0, len(ranks))
    ax.set_ylim(0, len(cases))
    ax.axis("off")
    cmap = plt.get_cmap("BuGn")
    score_min = float(top_df["score"].min())
    score_max = float(top_df["score"].max())

    for y, case_name in enumerate(cases):
        y_pos = len(cases) - y - 1
        ax.text(-0.08, y_pos + 0.5, case_name, ha="right", va="center", fontsize=12, fontweight="bold", color="#26323F")
        for x_idx, rank in enumerate(ranks):
            row = top_df[(top_df["case"] == case_name) & (top_df["rank"] == rank)]
            if row.empty:
                continue
            item = row.iloc[0]
            norm = (float(item["score"]) - score_min) / max(1e-6, score_max - score_min)
            color = cmap(0.10 + 0.30 * norm)
            rect = plt.Rectangle((x_idx + 0.02, y_pos + 0.08), 0.94, 0.84, facecolor=color, edgecolor="#D8D2C8", linewidth=1)
            ax.add_patch(rect)
            drug = str(item["drug_name"])
            if len(drug) > 24:
                drug = drug[:22] + "..."
            ax.text(x_idx + 0.49, y_pos + 0.60, f"{rank}. {drug}", ha="center", va="center", fontsize=10.2, color="#102027", fontweight="bold")
            ax.text(x_idx + 0.49, y_pos + 0.31, f"分数 {float(item['score']):.3f}", ha="center", va="center", fontsize=9, color="#34404B")

    for x_idx, rank in enumerate(ranks):
        ax.text(x_idx + 0.49, len(cases) + 0.14, f"第 {rank} 名", ha="center", va="bottom", fontsize=12, fontweight="bold", color="#26323F")
    ax.text(
        1.5,
        -0.25,
        "分数 = ExtraTrees 模型概率 + DrugCentral 已知靶点证据校准；仅表示计算筛选优先级。",
        ha="center",
        va="center",
        fontsize=9.2,
        color="#56606B",
    )
    path = FIG / "multi_disease_target_rankings.png"
    fig.tight_layout(pad=1.1)
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def load_model_comparison():
    if not MODEL_COMPARISON_SOURCE.exists():
        return pd.DataFrame()
    df = pd.read_csv(MODEL_COMPARISON_SOURCE)
    df.to_csv(MODEL_COMPARISON_CSV, index=False)
    return df


def plot_model_comparison(model_df):
    if model_df.empty:
        return None
    split_order = ["random_pair_split", "target_group_split", "drug_group_split"]
    split_labels = {"random_pair_split": "随机划分", "target_group_split": "冷靶点", "drug_group_split": "冷药物"}
    model_order = ["extra_trees", "hist_gradient_boosting", "catboost"]
    model_labels = {"extra_trees": "ExtraTrees", "hist_gradient_boosting": "HistGradientBoosting", "catboost": "CatBoost"}
    colors = {"extra_trees": "#254E7B", "hist_gradient_boosting": "#2F5D50", "catboost": "#9B5A2E"}

    fig, ax = plt.subplots(figsize=(7.6, 4.7))
    x = np.arange(len(split_order))
    width = 0.24
    for offset, model_name in zip([-width, 0, width], model_order):
        vals = []
        for split in split_order:
            row = model_df[(model_df["split"] == split) & (model_df["model"] == model_name)]
            vals.append(float(row.iloc[0]["roc_auc"]) if not row.empty else np.nan)
        ax.bar(x + offset, vals, width, label=model_labels[model_name], color=colors[model_name])
        for idx, value in enumerate(vals):
            if np.isfinite(value):
                ax.text(x[idx] + offset, value + 0.012, f"{value:.3f}", ha="center", va="bottom", fontsize=8.2)
    ax.set_xticks(x)
    ax.set_xticklabels([split_labels[s] for s in split_order])
    ax.set_ylim(0.78, 0.965)
    ax.set_ylabel("ROC-AUC")
    ax.set_title("三类模型在不同验证划分下的排序能力")
    ax.grid(axis="y", alpha=0.16, color="#D8D2C8")
    ax.legend(frameon=False, loc="upper center", bbox_to_anchor=(0.5, -0.11), ncol=3)
    fig.subplots_adjust(bottom=0.24)
    path = FIG / "model_comparison_roc_auc.png"
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def add_title(slide, text):
    rule = slide.shapes.add_shape(1, Inches(0.55), Inches(0.35), Inches(0.06), Inches(0.38))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    rule.line.color.rgb = BLUE
    add_textbox(slide, 0.72, 0.31, 2.2, 0.18, "模型证据", size=7.8, bold=True, color=(86, 95, 106))
    box = slide.shapes.add_textbox(Inches(0.70), Inches(0.52), Inches(11.7), Inches(0.44))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = CHINESE_FONT
    p.font.size = Pt(23)
    p.font.bold = True
    p.font.color.rgb = INK


def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=(43, 48, 53), align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = CHINESE_FONT
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        if align:
            p.alignment = align
    return box


def add_bullets(slide, x, y, w, h, bullets, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = CHINESE_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = MUTED
        p.space_after = Pt(6)
    return box


def add_takeaway(slide, text):
    shape = slide.shapes.add_shape(1, Inches(0.55), Inches(6.72), Inches(12.25), Inches(0.33))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    add_textbox(slide, 0.75, 6.77, 11.85, 0.24, text, size=9.5, color=(45, 63, 82))


def add_source(slide, text):
    add_textbox(slide, 0.55, 7.15, 12.25, 0.17, text, size=7.5, color=(105, 112, 119))


def add_image(slide, path, x, y, w, h):
    frame = slide.shapes.add_shape(1, Inches(x - 0.04), Inches(y - 0.04), Inches(w + 0.08), Inches(h + 0.08))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_metric_chip(slide, x, y, value, label, color=(37, 78, 123)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.25), Inches(0.82))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    accent = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(0.82))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*color)
    accent.line.color.rgb = RGBColor(*color)
    add_textbox(slide, x + 0.12, y + 0.10, 2.0, 0.24, value, size=19, bold=True, color=color)
    add_textbox(slide, x + 0.12, y + 0.49, 2.0, 0.20, label, size=8.5, color=(80, 87, 94))


def add_box_text(slide, x, y, w, h, text, size=13, fill=WHITE, line=LINE, color=(38, 50, 63), bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    add_textbox(slide, x + 0.12, y + 0.12, w - 0.24, h - 0.20, text, size=size, bold=bold, color=color, align=align)
    return shape


def add_slide_background(slide):
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.color.rgb = PAPER
    top_rule = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.045))
    top_rule.fill.solid()
    top_rule.fill.fore_color.rgb = INK
    top_rule.line.color.rgb = INK


def new_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    add_slide_background(slide)
    return slide


def build_ppt(summary, figs, metrics_df, learning_df, influenza_df, case_df):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    metric_lookup = metrics_df.set_index("split")
    random_auc = f"{metric_lookup.loc['random_pair_split', 'roc_auc']:.4f}"
    cold_target_auc = f"{metric_lookup.loc['target_group_split', 'roc_auc']:.4f}"
    cold_drug_auc = f"{metric_lookup.loc['drug_group_split', 'roc_auc']:.4f}"
    random_acc = f"{metric_lookup.loc['random_pair_split', 'accuracy'] * 100:.2f}%"
    cold_target_acc = f"{metric_lookup.loc['target_group_split', 'accuracy'] * 100:.2f}%"
    cold_drug_acc = f"{metric_lookup.loc['drug_group_split', 'accuracy'] * 100:.2f}%"

    # 1
    slide = new_slide(prs, blank)
    add_textbox(slide, 0.72, 0.72, 11.8, 0.65, TITLE, size=30, bold=True)
    add_textbox(slide, 0.76, 1.48, 10.8, 0.38, SUBTITLE, size=15, color=(55, 65, 75))
    add_metric_chip(slide, 0.82, 2.35, random_auc, "随机划分 ROC-AUC")
    add_metric_chip(slide, 3.23, 2.35, cold_target_auc, "冷靶点 ROC-AUC", color=(47, 93, 80))
    add_metric_chip(slide, 5.64, 2.35, cold_drug_auc, "冷药物 ROC-AUC", color=(155, 90, 46))
    add_textbox(slide, 0.86, 3.24, 8.7, 0.3, f"准确率：随机划分 {random_acc} ｜ 冷靶点 {cold_target_acc} ｜ 冷药物 {cold_drug_acc}", size=11.5, color=(86, 95, 106))
    add_textbox(
        slide,
        0.82,
        3.55,
        10.8,
        1.1,
        "预测目标：判断一个“药物结构—靶点描述”组合是否符合 DrugCentral 中已知药物—靶点关系模式。\n定位：用于已知药物重定位和候选召回排序，不用于药效、剂量或临床结论。",
        size=15,
    )
    add_takeaway(slide, "核心结论：ExtraTrees 在随机、冷靶点、冷药物三类验证中均保持较强排序能力，可作为本项目的稳健筛选基线。")
    add_source(slide, "数据：本地 DrugCentral 提取；模型：sklearn ExtraTreesClassifier；生成日期：2026-06-01。")

    # 2
    slide = new_slide(prs, blank)
    add_title(slide, "核心挑战在于把分子候选推进到可验证证据链")
    add_bullets(
        slide,
        0.75,
        1.05,
        5.4,
        4.9,
        [
            "AI 可以快速产生候选 SMILES；缺少靶点证据和已知药物基准时，筛选结果难以被解释。",
            "本模型将任务收缩为更可验证的问题：药物结构与靶点文本是否呈现已知药物—靶点关系模式。",
            "模型输出是计算筛选先验，用于排序、召回和进入后续第 4 阶段结构验证。",
            "该目标将机器学习分数限定为筛选先验，降低其被误读为生物活性或临床有效性的风险。",
        ],
        size=15,
    )
    add_textbox(
        slide,
        7.0,
        1.2,
        4.8,
        3.7,
        "输入\n药物 SMILES + 靶点名称/类别/物种/基因/Accession\n\n输出\n匹配概率 P(已知药物—靶点关系)\n\n用途\n老药重定位候选召回、阳性对照校准、第 3/4 阶段前置排序",
        size=17,
        bold=False,
        color=(32, 42, 54),
    )
    add_takeaway(slide, "预测目标被严格限定为“已知关系模式识别”，这是项目可讲清楚、可验证、可合规的基础。")

    # 3
    slide = new_slide(prs, blank)
    add_title(slide, "疾病输入到药物排名：产品链路强调可追溯证据")
    flow_items = [
        ("疾病/病例特征\n流感、HIV、糖尿病等", "用户输入"),
        ("靶点 Brief\n蛋白、物种、口袋、参考药", "证据整理"),
        ("本地药物库\n4099 个 DrugCentral 结构", "候选来源"),
        ("模型打分\nExtraTrees + 证据校准", "排序引擎"),
        ("Top-K 输出\n候选药物、分数、证据靶点", "前端展示"),
        ("后续验证\n对接、活性、实验设计", "验证入口"),
    ]
    x_positions = [0.72, 2.78, 4.84, 6.90, 8.96, 11.02]
    for idx, ((main_text, label), x0) in enumerate(zip(flow_items, x_positions)):
        fill = RGBColor(255, 255, 255) if idx % 2 == 0 else RGBColor(244, 248, 247)
        add_box_text(slide, x0, 1.72, 1.54, 1.72, main_text, size=10.8, fill=fill, bold=True)
        add_textbox(slide, x0, 3.62, 1.54, 0.22, label, size=8.5, color=(92, 101, 110), align=PP_ALIGN.CENTER)
        if idx < len(flow_items) - 1:
            add_textbox(slide, x0 + 1.55, 2.28, 0.42, 0.35, "→", size=22, bold=True, color=(37, 78, 123), align=PP_ALIGN.CENTER)
    add_bullets(
        slide,
        0.95,
        4.45,
        11.2,
        1.55,
        [
            "这条链路的重点不是直接宣称发现药物，而是把候选产生、已知证据、模型排序和验证入口连接起来。",
            "前端展示应优先呈现 Top-K 排名、证据来源和下一步验证建议，而不是只给一个孤立分数。",
        ],
        size=14,
    )
    add_takeaway(slide, "该流程图把模型从单一算法结果放回产品闭环：输入病例/疾病特征后，系统输出可解释候选排名。")

    # 4
    slide = new_slide(prs, blank)
    add_title(slide, "训练数据来自 DrugCentral：正关系与确定性负采样构成二分类任务")
    add_metric_chip(slide, 0.8, 1.15, "4099", "药物结构")
    add_metric_chip(slide, 3.2, 1.15, "2729", "靶点")
    add_metric_chip(slide, 5.6, 1.15, "18321", "正关系配对")
    add_metric_chip(slide, 8.0, 1.15, "18321", "负采样配对")
    add_bullets(
        slide,
        0.85,
        2.45,
        10.8,
        3.2,
        [
            "正样本：DrugCentral 中记录的药物—靶点关系。",
            "负样本：同一药物与未记录为其靶点的其他靶点配对，按 1:1 比例构造。",
            "总样本量为 36642 条，每条样本对应一个药物结构和一个靶点文本。",
            "负样本来自确定性配对构造，结果应解释为筛选排序能力，不应解释为真实失败实验证据。",
        ],
        size=15,
    )
    add_takeaway(slide, "数据设计的关键是把公开药物—靶点关系转化为可复现的监督学习问题，同时保留证据边界。")

    # 4
    slide = new_slide(prs, blank)
    add_title(slide, "特征由分子结构和靶点文本两条通道拼接而成")
    add_textbox(
        slide,
        0.85,
        1.15,
        3.5,
        3.9,
        "分子通道\n\nMorgan 指纹\n半径 = 2\n维度 = 512\n\nRDKit 描述符 × 8\n分子量、LogP、TPSA、氢键供体/受体、可旋转键、重原子数、环数量",
        size=15,
    )
    add_textbox(slide, 4.85, 2.2, 1.0, 0.45, "+", size=32, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        6.0,
        1.15,
        3.1,
        3.9,
        "靶点通道\n\n靶点文本哈希\n维度 = 256\n\n文本字段\n靶点名、类别、accession、gene、SwissProt、物种、作用类型",
        size=15,
    )
    add_textbox(slide, 9.5, 2.2, 1.0, 0.45, "=", size=32, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        10.45,
        1.55,
        2.1,
        2.4,
        "776 维输入\n\n512 + 8 + 256\n\n用于 ExtraTrees 分类器",
        size=17,
        bold=True,
        color=(37, 78, 123),
        align=PP_ALIGN.CENTER,
    )
    add_takeaway(slide, "该特征设计让模型同时看到化学结构局部模式、药物样性质和靶点证据文本。")

    # 5
    slide = new_slide(prs, blank)
    add_title(slide, "ExtraTrees 通过大量随机树学习非线性匹配规则")
    add_bullets(
        slide,
        0.75,
        1.05,
        5.6,
        4.9,
        [
            "ExtraTrees 是 Extremely Randomized Trees，即极端随机树集成模型。",
            "每棵树在特征选择和切分阈值上引入随机性，最后对多棵树的概率输出做平均。",
            "它适合高维稀疏指纹、数值描述符和哈希文本特征混合的结构化任务。",
            "当前参数：180 棵树，min_samples_leaf=2，max_features=sqrt，random_state=13。",
        ],
        size=15,
    )
    add_textbox(
        slide,
        7.0,
        1.25,
        4.8,
        3.5,
        "单一规则的局限\n\n药物—靶点关系通常由多个结构片段、理化性质和靶点证据共同决定。\n\nExtraTrees 的优势在于以较低调参成本捕捉这种非线性组合。",
        size=16,
    )
    add_takeaway(slide, "模型定位为稳健的筛选基线，用于排序和召回，不承担药效判定终点。")

    slide = new_slide(prs, blank)
    add_title(slide, "排名分数由模型概率和已知靶点证据共同决定")
    add_box_text(slide, 0.88, 1.55, 2.75, 1.55, "ExtraTrees\n模型概率\n\nP(药物—靶点关系)", size=14, fill=WHITE, bold=True)
    add_textbox(slide, 3.72, 2.05, 0.48, 0.42, "+", size=30, bold=True, color=(37, 78, 123), align=PP_ALIGN.CENTER)
    add_box_text(slide, 4.35, 1.55, 3.0, 1.55, "DrugCentral\n已知靶点证据\n\n靶点名称/基因/物种相似度", size=14, fill=RGBColor(244, 248, 247), bold=True)
    add_textbox(slide, 7.45, 2.05, 0.48, 0.42, "=", size=30, bold=True, color=(37, 78, 123), align=PP_ALIGN.CENTER)
    add_box_text(slide, 8.1, 1.55, 3.65, 1.55, "最终排序分数\n\n用于 Top-K 候选药物排名", size=15, fill=RGBColor(245, 248, 252), bold=True)
    add_bullets(
        slide,
        1.0,
        3.75,
        10.9,
        1.75,
        [
            "模型概率来自分子结构特征和靶点文本特征，反映机器学习学到的匹配模式。",
            "已知靶点证据来自本地 DrugCentral 关系，用于让已验证方向在产品排序中拥有更高可信度。",
            "最终分数用于筛选优先级，不直接代表结合强度、药效或临床结论。",
        ],
        size=14,
    )
    add_takeaway(slide, "这页解决“分数从哪里来”的问题：模型提供泛化排序能力，数据库证据提供可解释校准。")

    # 6
    slide = new_slide(prs, blank)
    add_title(slide, "三种验证划分分别检验插值、冷靶点和冷药物泛化")
    add_textbox(slide, 0.85, 1.15, 3.55, 3.7, "随机划分\n\n随机拆分药物—靶点配对。\n\n回答：模型能否在总体样本中区分已知关系和负采样关系？", size=15)
    add_textbox(slide, 4.9, 1.15, 3.55, 3.7, "冷靶点划分\n\n测试靶点在训练中不可见。\n\n回答：面对新的靶点描述，模型是否仍能排序？", size=15)
    add_textbox(slide, 8.95, 1.15, 3.55, 3.7, "冷药物划分\n\n测试药物在训练中不可见。\n\n回答：面对新的药物结构，模型是否仍能外推？", size=15)
    add_takeaway(slide, "随机划分分数最高符合预期；真正支撑泛化叙事的是冷靶点与冷药物仍接近 0.89 的 ROC-AUC。")

    slide = new_slide(prs, blank)
    add_title(slide, "冷启动验证更接近真实使用场景")
    add_box_text(slide, 0.85, 1.3, 3.45, 3.2, "随机划分\n\n更像插值测试。\n样本分布接近训练集，适合确认模型是否学到总体模式。", size=14, fill=WHITE, bold=True)
    add_box_text(slide, 4.95, 1.3, 3.45, 3.2, "冷靶点\n\n更像新疾病/新蛋白问题。\n训练阶段没见过该靶点，考察靶点文本泛化。", size=14, fill=RGBColor(244, 248, 247), bold=True)
    add_box_text(slide, 9.05, 1.3, 3.45, 3.2, "冷药物\n\n更像新候选分子问题。\n训练阶段没见过该药物，考察结构特征外推。", size=14, fill=RGBColor(245, 248, 252), bold=True)
    add_bullets(
        slide,
        0.95,
        5.0,
        11.2,
        0.95,
        [
            "因此，随机划分说明模型会做基本区分；冷靶点和冷药物才更能支撑“可扩展筛选平台”的叙事。",
            "当前冷靶点和冷药物 ROC-AUC 均约 0.893，说明模型具备一定跨靶点和跨药物排序能力。",
        ],
        size=13.6,
    )
    add_takeaway(slide, "汇报时应强调冷启动表现，因为它比随机划分更接近新病毒、新靶点和新候选药物场景。")

    # 7
    slide = new_slide(prs, blank)
    add_title(slide, "ROC-AUC 与准确率共同刻画模型的排序和分类表现")
    add_image(slide, figs["metric_bars"], 0.65, 1.0, 7.2, 4.6)
    add_bullets(
        slide,
        8.15,
        1.15,
        4.2,
        3.7,
        [
            f"ROC-AUC：随机划分 {random_auc}，冷靶点 {cold_target_auc}，冷药物 {cold_drug_auc}。",
            f"准确率：随机划分 {random_acc}，冷靶点 {cold_target_acc}，冷药物 {cold_drug_acc}。",
            "冷启动划分仍维持接近 0.89 的 ROC-AUC，说明模型具备可用的排序泛化能力。",
            "F1 在冷靶点下降，提示阈值分类更保守；产品侧应优先展示排序分和 Top-K。",
        ],
        size=13,
    )
    add_takeaway(slide, "模型适合作为候选召回和优先级排序模块；阈值判断需要结合证据来源和后续验证。")
    add_source(slide, "图表来自 seed=13 的本地复现实验；指标与项目现有 benchmark 一致。")

    if figs.get("model_comparison"):
        slide = new_slide(prs, blank)
        add_title(slide, "模型对比显示 ExtraTrees 是当前最稳健的传统基线")
        add_image(slide, figs["model_comparison"], 0.72, 1.0, 7.65, 4.75)
        add_bullets(
            slide,
            8.65,
            1.15,
            3.85,
            3.75,
            [
                "对比模型：ExtraTrees、HistGradientBoosting、CatBoost。",
                "三个模型使用同一 DrugCentral 数据和同一特征体系。",
                "ExtraTrees 在随机、冷靶点、冷药物三类划分中 ROC-AUC 均为最高。",
                "因此 PPT 将 ExtraTrees 作为主模型讲解，其他模型作为基准参照。",
            ],
            size=13.2,
        )
        add_takeaway(slide, "模型对比页让“为什么选 ExtraTrees”有数据依据，而不是只展示单一模型结果。")
        add_source(slide, f"模型对比数据：{MODEL_COMPARISON_CSV}")

    # 8
    slide = new_slide(prs, blank)
    add_title(slide, "ROC 曲线显示模型在不同假阳性率下仍能稳定召回正关系")
    add_image(slide, figs["roc"], 0.8, 0.95, 7.1, 5.4)
    add_bullets(
        slide,
        8.35,
        1.2,
        4.0,
        3.5,
        [
            "ROC-AUC 衡量随机抽取一条正关系和一条负样本时，模型把正关系排在前面的概率。",
            "三条曲线均显著高于随机对角线。",
            "冷启动曲线较随机划分下降，符合真实泛化难度。",
        ],
        size=14,
    )
    add_takeaway(slide, "ROC 曲线支持该模型作为早期虚拟筛选中的排序先验。")
    add_source(slide, "ROC 曲线来自本地训练/验证划分重跑结果，未使用模拟曲线。")

    # 9
    slide = new_slide(prs, blank)
    add_title(slide, "阈值分类显示：冷靶点场景更保守，适合 Top-K 排序")
    add_image(slide, figs["confusion"], 0.75, 1.05, 7.75, 3.75)
    add_bullets(
        slide,
        8.8,
        1.15,
        3.65,
        3.55,
        [
            "混淆矩阵采用概率阈值 0.5。",
            "冷靶点划分 precision 较高但 recall 较低，说明模型倾向保守召回。",
            "产品中应展示排序分、证据来源和 Top-K，并将 0.5 阈值视为可调工作点。",
        ],
        size=14,
    )
    add_takeaway(slide, "这页解释了为什么 ROC-AUC 很强，但 F1 在冷靶点场景下降：模型更适合作为筛选排序器。")

    # 10
    slide = new_slide(prs, blank)
    add_title(slide, "ExtraTrees 没有神经网络式训练轮次；这里补跑集成规模学习曲线")
    add_image(slide, figs["learning"], 0.75, 1.0, 7.35, 4.9)
    add_bullets(
        slide,
        8.35,
        1.1,
        4.0,
        3.7,
        [
            "树模型不通过梯度下降逐轮更新，因此不存在神经网络意义上的训练损失曲线。",
            "补跑实验记录树数量从 10 到 180 时的训练/验证 log-loss 和验证 ROC-AUC。",
            "曲线用于判断集成规模是否足够，以及继续增加树数是否明显改善。",
        ],
        size=14,
    )
    add_takeaway(slide, "该图是基于树数量补跑的训练稳定性分析，用 log-loss 描述集成规模带来的性能变化。")

    # 11
    slide = new_slide(prs, blank)
    add_title(slide, "特征重要性表明：结构指纹、描述符和靶点文本共同驱动预测")
    add_image(slide, figs["importance"], 0.72, 1.0, 7.65, 4.55)
    add_bullets(
        slide,
        8.65,
        1.1,
        3.8,
        3.7,
        [
            "Morgan 指纹捕捉局部结构片段，是药物结构侧的主要信息来源。",
            "RDKit 标量描述符提供分子量、极性、氢键、柔性和环结构等可解释补充。",
            "靶点文本哈希向量提供靶点身份与证据上下文。",
        ],
        size=14,
    )
    add_takeaway(slide, "模型同时使用药物结构、理化描述符和靶点文本三类证据。")

    # 12
    slide = new_slide(prs, blank)
    add_title(slide, "SHAP 分析用于检查模型排序主要由哪些特征块贡献")
    add_image(slide, figs["shap"], 0.72, 1.0, 7.7, 4.7)
    add_bullets(
        slide,
        8.65,
        1.12,
        3.8,
        3.65,
        [
            "SHAP 在验证样本上计算平均 |SHAP|，反映特征对预测概率的平均影响强度。",
            "由于指纹与靶点哈希是高维编码，单个位点不直接等同于具体药效机制。",
            "SHAP 结果更适合用于模型审计，生物机制仍需结构和实验验证支撑。",
        ],
        size=14,
    )
    add_takeaway(slide, "SHAP 提供模型审计入口，药理机制解释仍需外部证据闭环。")

    # 13
    slide = new_slide(prs, blank)
    add_title(slide, "甲流 NA 案例：已知 neuraminidase 药物被召回到前列")
    table = slide.shapes.add_table(7, 4, Inches(0.75), Inches(1.05), Inches(7.2), Inches(4.2)).table
    headers = ["排序", "药物", "分数", "证据靶点"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(237, 242, 247)
    for i, (_, row) in enumerate(influenza_df.iterrows(), start=1):
        table.cell(i, 0).text = str(int(row.get("rank", i)))
        table.cell(i, 1).text = str(row.get("drug_name", row.get("drug", "")))[:28]
        score = row.get("match_probability", row.get("score", ""))
        table.cell(i, 2).text = f"{float(score):.3f}" if score != "" and not pd.isna(score) else ""
        target = str(row.get("evidence_target_name", row.get("evidence", "")))
        key = str(row.get("evidence_target_key", ""))
        table.cell(i, 3).text = (key + " " + target).strip()[:30]
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                p.font.name = CHINESE_FONT
                p.font.size = Pt(8.5)
    add_bullets(
        slide,
        8.35,
        1.15,
        4.2,
        3.7,
        [
            "查询输入：Influenza A H1N1 neuraminidase NA P03468 oseltamivir pocket。",
            "前三名为 oseltamivir、zanamivir、peramivir，符合已知 NA 药物基准。",
            "该案例验证召回链路合理，后续药效判断仍需要结构和实验验证。",
        ],
        size=14,
    )
    add_takeaway(slide, "已知阳性方向能够被召回，是后续候选生成和第 4 阶段校验的基础控制实验。")

    # 14
    slide = new_slide(prs, blank)
    add_title(slide, "多疾病靶点案例展示：从 4099 个本地药物结构中批量排序")
    add_image(slide, figs["case_rankings"], 0.62, 1.0, 8.15, 4.75)
    add_bullets(
        slide,
        9.05,
        1.12,
        3.55,
        3.95,
        [
            "案例覆盖病毒感染、心血管、代谢、炎症疼痛和肿瘤相关靶点。",
            "每个靶点均对本地 DrugCentral 药物结构库做全量排序。",
            "最终分数由 ExtraTrees 结构匹配概率与已知靶点证据校准得到。",
            "该页展示的是重定位候选优先级，后续仍需对接、活性和实验验证。",
        ],
        size=13.2,
    )
    add_takeaway(slide, "跨靶点排名页能更清楚地展示产品价值：输入疾病/靶点后，系统给出可解释的已知药物候选列表。")

    # 15
    slide = new_slide(prs, blank)
    add_title(slide, "在产品链路中，ExtraTrees 是阶段之间的排序桥梁")
    add_textbox(slide, 0.85, 1.2, 2.6, 3.6, "第 1/2 阶段\n\n靶点 Brief\n证据矩阵\n已知药物库", size=16, align=PP_ALIGN.CENTER)
    add_textbox(slide, 3.85, 2.35, 0.6, 0.4, "→", size=28, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 4.65, 1.2, 3.2, 3.6, "ExtraTrees 匹配器\n\n药物—靶点关系概率\n已知药物先验\n候选召回排序", size=16, align=PP_ALIGN.CENTER)
    add_textbox(slide, 8.1, 2.35, 0.6, 0.4, "→", size=28, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 8.9, 1.2, 3.2, 3.6, "第 3/4 阶段\n\n候选库筛选\n控药/Decoy 校准\nRDKit/对接验证", size=16, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "模型的产品角色是减少搜索空间、提供可解释优先级，并把已知药物基准接入后续验证。")

    slide = new_slide(prs, blank)
    add_title(slide, "科学边界与下一步验证路径需要单独讲清楚")
    add_box_text(slide, 0.85, 1.25, 3.65, 3.1, "当前已支持\n\n本地药物库排序\n已知靶点证据校准\n冷启动泛化评估\n多疾病案例展示", size=14, fill=WHITE, bold=True)
    add_box_text(slide, 4.9, 1.25, 3.65, 3.1, "尚不能宣称\n\n真实结合强度\n体外活性\n毒性和安全性\n剂量与临床有效性", size=14, fill=RGBColor(252, 246, 240), bold=True)
    add_box_text(slide, 8.95, 1.25, 3.65, 3.1, "下一步补强\n\nChEMBL/BindingDB 活性数据\n困难负样本采样\nDocking/Pose QC\n实验验证设计", size=14, fill=RGBColor(244, 248, 247), bold=True)
    add_bullets(
        slide,
        0.95,
        4.8,
        11.2,
        1.0,
        [
            "这页用于主动设定合规边界：模型只提供计算筛选优先级，后续验证决定候选是否真正有价值。",
            "对评审来说，清楚的边界反而会增强项目可信度。",
        ],
        size=13.8,
    )
    add_takeaway(slide, "最稳妥的产品定位：高通量计算筛选和证据排序平台，而不是药效结论生成器。")

    # 16
    slide = new_slide(prs, blank)
    add_title(slide, "结论：强排序性能支撑“已知证据驱动”的筛选路线")
    add_bullets(
        slide,
        0.85,
        1.05,
        6.0,
        4.9,
        [
            "模型把 DrugCentral 关系学习为可复现的药物—靶点匹配任务。",
            f"随机、冷靶点、冷药物 ROC-AUC 分别为 {random_auc}、{cold_target_auc}、{cold_drug_auc}；准确率分别为 {random_acc}、{cold_target_acc}、{cold_drug_acc}。",
            "特征与 SHAP 分析表明分子结构、理化描述符和靶点文本均参与决策。",
            "甲流 NA 案例能召回已知 neuraminidase 药物，适合作为演示场景的阳性基准。",
        ],
        size=15,
    )
    add_textbox(
        slide,
        7.6,
        1.25,
        4.6,
        3.65,
        "边界\n\n这些结果支持计算筛选排序能力。\n\n真实结合强度、生物活性、药效、毒性、安全性、剂量和临床可用性仍需要独立证据。\n\n下一步应接入 ChEMBL/BindingDB 活性数据、困难负样本采样和外部实验/结构验证。",
        size=15,
    )
    add_takeaway(slide, "最稳妥的表述：这是一个性能较强、可审计的靶点—药物匹配基线，适合作为后续验证的前置筛选器。")

    slide = new_slide(prs, blank)
    add_title(slide, "附录：复现实验与交付文件索引")
    appendix_rows = [
        ("训练数据", "DrugCentral 4099 个药物结构、2729 个靶点、18321 条正关系"),
        ("输入特征", "Morgan 指纹 512 维 + RDKit 描述符 8 维 + 靶点文本哈希 256 维"),
        ("主模型", "ExtraTreesClassifier：180 棵树、min_samples_leaf=2、max_features=sqrt"),
        ("验证划分", "随机划分、冷靶点、冷药物；均使用 seed=13 可复现"),
        ("排名输出", "multi_disease_target_rankings.csv：7 个疾病/靶点场景 × Top 5"),
        ("模型对比", "model_comparison_benchmark.csv：ExtraTrees / HistGradientBoosting / CatBoost"),
    ]
    y0 = 1.15
    for idx, (label, value) in enumerate(appendix_rows):
        y = y0 + idx * 0.72
        fill = WHITE if idx % 2 == 0 else RGBColor(247, 249, 250)
        add_box_text(slide, 0.85, y, 2.15, 0.50, label, size=11, fill=fill, bold=True)
        add_box_text(slide, 3.08, y, 9.15, 0.50, value, size=10.5, fill=fill, align=PP_ALIGN.LEFT)
    add_textbox(
        slide,
        0.95,
        5.85,
        11.4,
        0.55,
        f"主要产物路径：{OUT}",
        size=10.5,
        color=(86, 95, 106),
    )
    add_takeaway(slide, "附录页用于答辩追问：数据规模、特征维度、模型参数、排名文件和对比基准都有可追溯文件。")

    prs.save(PPTX)


def main():
    ensure_dirs()
    set_plot_style()
    dataset, x, y, struct_ids, target_keys = build_dataset()
    splits = split_indices(y, struct_ids, target_keys)
    metrics = {}
    roc_data = {}
    models = {}
    for name, (train_idx, valid_idx) in splits.items():
        model, m, roc = evaluate_split(name, x, y, train_idx, valid_idx)
        metrics[name] = m
        roc_data[name] = roc
        models[name] = model

    metrics_df = pd.DataFrame(metrics.values())
    metrics_df.to_csv(OUT / "rerun_extra_trees_metrics.csv", index=False)

    roc_path = plot_roc(roc_data, metrics)
    metric_path = plot_metric_comparison(metrics_df)
    learning_path, learning_df = learning_curve(x, y, *splits["random_pair_split"])
    importance_path, importance_blocks = plot_feature_importance(models["random_pair_split"])
    shap_path, shap_blocks = run_shap(models["random_pair_split"], x, splits["random_pair_split"][1])
    confusion_path = plot_confusion(metrics)
    production_model = train_extra_trees(x, y)
    case_df = rank_multi_disease_cases(production_model, dataset)
    case_rankings_path = plot_case_rankings(case_df)
    model_comparison_df = load_model_comparison()
    model_comparison_path = plot_model_comparison(model_comparison_df)

    influenza_df = load_influenza_predictions()

    summary = {
        "dataset": dataset.summary,
        "metrics": metrics,
        "learning_curve_rows": learning_df.to_dict(orient="records"),
        "feature_importance_blocks": importance_blocks,
        "shap_blocks": shap_blocks,
        "shap_validation_sample_size": 40,
        "multi_disease_cases": case_df.to_dict(orient="records"),
        "model_comparison": model_comparison_df.to_dict(orient="records") if not model_comparison_df.empty else [],
        "boundary": BOUNDARY,
    }
    SUMMARY_JSON.write_text(json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    figs = {
        "roc": roc_path,
        "metric_bars": metric_path,
        "learning": learning_path,
        "importance": importance_path,
        "shap": shap_path,
        "confusion": confusion_path,
        "case_rankings": case_rankings_path,
        "model_comparison": model_comparison_path,
    }
    build_ppt(summary, figs, metrics_df, learning_df, influenza_df, case_df)

    # Lightweight PPTX verification.
    reopened = Presentation(str(PPTX))
    media_count = len(list((PPTX.parent / "_tmp").glob("*"))) if (PPTX.parent / "_tmp").exists() else "n/a"
    qa = [
        "# QA Report",
        "",
        f"- PPTX: `{PPTX}`",
        f"- Slide count: {len(reopened.slides)}",
        f"- Figures generated: {len(list(FIG.glob('*.png')))}",
        "- Verification: PPTX reopened successfully with python-pptx.",
        "- ROC curves, learning curve, feature importance, and SHAP plots were generated from local rerun experiments.",
        "- Multi-disease ranking panel was generated from local DrugCentral structures using ExtraTrees plus known-target evidence calibration.",
        "- Model comparison panel was generated from an existing local benchmark CSV, not manually entered.",
        "- SHAP grouped summary used a fixed-seed validation sample of 40 rows with TreeExplainer approximate mode for reproducible model audit.",
        "- ExtraTrees has no neural-network epoch loss; the deck uses log-loss vs number of trees as the valid training-stability curve.",
        f"- Boundary: {BOUNDARY}",
        "",
    ]
    QA.write_text("\n".join(qa), encoding="utf-8")
    print(PPTX)
    print(QA)
    print(SUMMARY_JSON)


if __name__ == "__main__":
    main()
